from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE
import os

prs = Presentation()
prs.slide_width = Inches(10)
prs.slide_height = Inches(5.625) # 16:9 aspect ratio is typical for modern presentations

# 3 images paths:
cover_img = r"C:\Users\126558_kh\.gemini\antigravity\brain\bbd59b1c-99dc-45db-88a9-d90838e87d0a\ppt_cover_dashboard_1773816315594.png"
sim_img = r"C:\Users\126558_kh\.gemini\antigravity\brain\bbd59b1c-99dc-45db-88a9-d90838e87d0a\ppt_map_simulation_1773816330291.png"
analytics_img = r"C:\Users\126558_kh\.gemini\antigravity\brain\bbd59b1c-99dc-45db-88a9-d90838e87d0a\ppt_ai_analytics_1773816345056.png"

def add_image_slide(title, subtitle, bullets, img_path):
    slide_layout = prs.slide_layouts[6] # Blank
    slide = prs.slides.add_slide(slide_layout)
    
    # Set background
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(15, 23, 42) # #0f172a
    
    # Add image on the right
    if os.path.exists(img_path):
        slide.shapes.add_picture(img_path, Inches(5.3), Inches(0.5), width=Inches(4.4), height=Inches(4.6))
        
    # Title
    txBox = slide.shapes.add_textbox(Inches(0.4), Inches(0.4), Inches(4.5), Inches(1))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.text = title
    p.font.bold = True
    p.font.size = Pt(24)
    p.font.color.rgb = RGBColor(56, 189, 248) # Neon blue
    
    # Subtitle
    if subtitle:
        p2 = tf.add_paragraph()
        p2.text = subtitle
        p2.font.size = Pt(14)
        p2.font.color.rgb = RGBColor(148, 163, 184)
        p2.space_before = Pt(8)
    
    # Body
    bodyBox = slide.shapes.add_textbox(Inches(0.4), Inches(1.6), Inches(4.6), Inches(3.5))
    tf_body = bodyBox.text_frame
    tf_body.word_wrap = True
    for bullet in bullets:
        p = tf_body.add_paragraph()
        p.text = "• " + bullet
        p.font.size = Pt(12)
        p.font.color.rgb = RGBColor(226, 232, 240)
        p.space_after = Pt(12)
        
# Cover Title Slide
slide = prs.slides.add_slide(prs.slide_layouts[6])
slide.background.fill.solid()
slide.background.fill.fore_color.rgb = RGBColor(15, 23, 42)
if os.path.exists(cover_img):
    slide.shapes.add_picture(cover_img, Inches(0), Inches(0), width=Inches(10), height=Inches(5.625))
    
    # Overlay transparent dark box
    shape = slide.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(0), Inches(1.5), Inches(10), Inches(2.2))
    shape.fill.solid()
    shape.fill.fore_color.rgb = RGBColor(0, 0, 0)
    shape.fill.transparency = 0.4
    shape.line.fill.background()

txBox = slide.shapes.add_textbox(Inches(0.5), Inches(1.8), Inches(9), Inches(1.8))
tf = txBox.text_frame
p = tf.paragraphs[0]
p.text = "AquaMoniq: AI 기반 하천 수질 모니터링 플랫폼"
p.font.bold = True
p.font.size = Pt(32)
p.font.color.rgb = RGBColor(255, 255, 255)
p.alignment = PP_ALIGN.CENTER

p2 = tf.add_paragraph()
p2.text = "현실 세계의 불확실성을 AI 예측 기술로 극복하다"
p2.font.size = Pt(18)
p2.font.color.rgb = RGBColor(56, 189, 248)
p2.alignment = PP_ALIGN.CENTER
p2.space_before = Pt(10)

# Slide 2: Challenge
add_image_slide(
    "1. 비점오염 모니터링의 한계", 
    "변덕스러운 강우, 그리고 놓쳐버린 골든타임", 
    [
        "기상 예보의 한계: 강우 예측(시간 및 강수량)과 실제 현장 발생 간의 빈번한 불일치",
        "추적의 패러독스: 유역에 내린 빗물이 여러 지류를 모아 하우스 형태의 조사 지점까지 오염원을 이끌고 도달하는 정확한 시점을 모름",
        "비효율성의 극대화: 부적절한 타이밍에 출동하여 조사를 실패하거나 인력을 낭비하는 악순환 발생"
    ], 
    analytics_img
)

# Slide 3: Solution
add_image_slide(
    "2. 지능형 강우 예측 엔진 도입",
    "KMA(기상청) 초단기 강수 예측 시스템 연동",
    [
        "기상청(KMA) 실황 및 초단기 예측망(VSRF) 가동",
        "자체 수학적 시뮬레이션 모델과 결합하여 예측 신뢰도 향상",
        "강우 유입 정확도 향상으로 헛걸음 방지",
        "강우 레이더 망 시각화를 통한 직관적 기상 실황 판단 지원"
    ],
    cover_img
)

# Slide 4: Path
add_image_slide(
    "3. 지능형 유역 공간 시뮬레이터",
    "오염물질의 흐름을 초단위로 낱낱이 파헤치다",
    [
        "정밀한 수계 지형 분석: 최적 채수 지점(Outlet)을 중심으로 유역을 눈물방울 형태로 재구성 및 시각화",
        "Pathfinder Routing AI: 빗물 속의 오염원 입자(Particle)가 수많은 강줄기를 관통하며 본류로 합쳐지는 궤적을 3D 애니메이션 형태로 제공",
        "정확한 도달 시간(ETA) 산출: 유출 경로 최적해 알고리즘 기반으로 도달 시간을 분 단위까지 예측, 조사관에게 최적의 채수 '골든타임' 통보"
    ],
    sim_img
)

# Slide 5: Realtime
add_image_slide(
    "4. AI 수질 오염 예측 모델",
    "센서가 없어도 실시간 데이터를 추론하다",
    [
        "빅데이터 딥러닝: 과거 다년간의 강수량-수질(BOD, T-P) 매핑 데이터를 학습한 딥러닝 알고리즘 활용",
        "실시간 오염 부하량 추정: 측정 장비가 직접 수집하지 않아도, 현재 유입되는 빗물의 양을 역산하여 시간대별 오염 물질 누적량 자동 시각화",
        "원클릭 통합 시연(Simulation System): 클릭 한 번으로 전체 인과과정 통합 시뮬레이션 지원"
    ],
    analytics_img
)

# Slide 6: Vision
add_image_slide(
    "5. 결론 및 미래 비전",
    "수질 관리 패러다임의 명백한 게임 체인저",
    [
        "경제적 효율성 극대화: 헛걸음률 0%, 불필요한 현장 출동 방지",
        "모니터링 정확도 및 성공률 획기적 향상",
        "미래 확장 교두보 확보: 사물인터넷(IoT), CCTV 스마트 수위 측정 망과 연동한 완전 무인 수질 관제 시스템 구축 예정"
    ],
    cover_img
)

prs.save(r"d:\(ai관련)\비점오염_AI_예측프로젝트\AquaMoniq_AI_Presentation.pptx")
print("Presentation generated successfully!")
